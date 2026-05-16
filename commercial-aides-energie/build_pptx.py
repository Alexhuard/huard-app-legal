"""
Generation du PowerPoint commercial - Aides Energie Groupe HUARD 2026
Support de presentation pour COMEX / equipe commerciale (16 slides).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = "/home/user/huard-app-legal/commercial-aides-energie/Aides-Primes-Energie-HUARD-2026.pptx"

NAVY = RGBColor(0x1F, 0x2F, 0x4D)
BLUE = RGBColor(0x5D, 0x81, 0xA6)
LIGHT = RGBColor(0xE8, 0xEE, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x5A, 0x64, 0x78)
GREEN = RGBColor(0x2A, 0x9D, 0x8F)
ORANGE = RGBColor(0xF4, 0xA2, 0x61)
RED = RGBColor(0xE7, 0x6F, 0x51)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background() if line is None else None
    if line is not None:
        shape.line.color.rgb = line
    return shape


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=NAVY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_header(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.9), NAVY)
    add_rect(slide, Inches(0), Inches(0.85), prs.slide_width, Inches(0.06), BLUE)
    add_text(slide, Inches(0.4), Inches(0.18), Inches(12), Inches(0.6), title, size=24, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.4), Inches(0.55), Inches(12), Inches(0.3), subtitle, size=11, color=LIGHT)


def add_footer(slide, page_num, total=22):
    add_text(slide, Inches(0.4), Inches(7.1), Inches(6), Inches(0.3),
             "Groupe HUARD - Aides & Primes Energie 2026 - Document interne", size=9, color=GREY)
    add_text(slide, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
             f"{page_num}/{total}", size=9, color=GREY, align=PP_ALIGN.RIGHT)


# =========================================================================
# SLIDE 1 - COVER
# =========================================================================
s = add_blank()
add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, NAVY)
add_rect(s, Inches(0), Inches(3.4), prs.slide_width, Inches(0.08), ORANGE)

add_text(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.6),
         "GROUPE HUARD", size=20, bold=True, color=ORANGE)
add_text(s, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.2),
         "Aides & Primes Energie 2026", size=44, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(3.7), Inches(11.5), Inches(0.6),
         "Levier commercial - Recensement des dispositifs mobilisables pour nos clients",
         size=18, color=LIGHT)
add_text(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.4),
         "Electricite - CVC - Telecoms - Informatique - Maintenance", size=14, color=LIGHT)
add_text(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.4),
         "Direction commerciale - 16 mai 2026", size=12, color=ORANGE)
add_text(s, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4),
         "Document interne - non contractuel", size=10, color=GREY)

# =========================================================================
# SLIDE 2 - POURQUOI CE DOCUMENT
# =========================================================================
s = add_blank()
add_header(s, "1. Pourquoi ce document", "Un outil pour transformer les aides en avantage commercial")

cards = [
    ("Constat", ORANGE,
     "Nos clients connaissent mal le paysage des aides. Resultat : ils repoussent les chantiers ou les paient en cash, alors que 25 a 70% peut etre finance."),
    ("Marche", GREEN,
     "Decret Tertiaire, Decret BACS, Loi APER, prime solaire renforcee 2026 : la reglementation cree un marche captif sur tous nos metiers."),
    ("Notre opportunite", BLUE,
     "En maitrisant 37 dispositifs, HUARD devient le partenaire qui simplifie la vie du client : on chiffre, on monte les aides, on realise, on entretient."),
]

x = Inches(0.5)
y = Inches(1.5)
w = Inches(4.0)
h = Inches(5.0)
gap = Inches(0.2)

for i, (title, color, body) in enumerate(cards):
    sx = x + (w + gap) * i
    add_rect(s, sx, y, w, Inches(0.6), color)
    add_text(s, sx + Inches(0.2), y + Inches(0.12), w - Inches(0.4), Inches(0.4),
             title, size=18, bold=True, color=WHITE)
    add_rect(s, sx, y + Inches(0.6), w, h - Inches(0.6), LIGHT)
    add_text(s, sx + Inches(0.3), y + Inches(0.9), w - Inches(0.6), h - Inches(1.0),
             body, size=13, color=NAVY)

add_text(s, Inches(0.5), Inches(6.7), Inches(12), Inches(0.4),
         "OBJECTIF : faire des aides un levier de signature (delais raccourcis, conversion devis amelioree, ticket moyen plus eleve).",
         size=12, bold=True, color=NAVY)
add_footer(s, 2)

# =========================================================================
# SLIDE 3 - CARTOGRAPHIE
# =========================================================================
s = add_blank()
add_header(s, "2. Cartographie des dispositifs", "37 aides cartographiees - 9 categories")

categories = [
    ("CEE Tertiaire", "9 fiches", BLUE),
    ("CEE Industrie", "3 fiches", BLUE),
    ("Coup de pouce CEE", "3 dispositifs", ORANGE),
    ("Photovoltaique", "5 dispositifs", GREEN),
    ("ADEME / France 2030", "5 dispositifs", GREEN),
    ("Bpifrance / Financement", "2 prets", BLUE),
    ("IRVE - bornes recharge", "2 dispositifs", ORANGE),
    ("Obligations creatrices de marche", "2 dispositifs", RED),
    ("Region IDF + collectivites", "4 dispositifs", BLUE),
    ("Residentiel / Copropriete", "5 dispositifs", GREEN),
    ("Programmes specialises", "3 dispositifs", GREY),
]

cols = 4
cw = Inches(2.95)
ch = Inches(1.2)
gap_x = Inches(0.15)
gap_y = Inches(0.2)
sx0 = Inches(0.5)
sy0 = Inches(1.4)

for i, (cat, count, color) in enumerate(categories):
    col = i % cols
    row = i // cols
    sx = sx0 + (cw + gap_x) * col
    sy = sy0 + (ch + gap_y) * row
    add_rect(s, sx, sy, cw, ch, color)
    add_text(s, sx + Inches(0.2), sy + Inches(0.15), cw - Inches(0.4), Inches(0.5),
             cat, size=13, bold=True, color=WHITE)
    add_text(s, sx + Inches(0.2), sy + Inches(0.65), cw - Inches(0.4), Inches(0.5),
             count, size=11, color=LIGHT)

add_text(s, Inches(0.5), Inches(6.0), Inches(12), Inches(1.0),
         "9 categories - 37 dispositifs - tous detailles dans le classeur Excel d'accompagnement (onglet 01 'Tableau complet').",
         size=12, color=NAVY)
add_text(s, Inches(0.5), Inches(6.4), Inches(12), Inches(0.6),
         "Code couleur : BLEU = CEE & financement / ORANGE = bonifications & IRVE / VERT = ADEME, EnR et copro / ROUGE = obligations reglementaires.",
         size=11, color=GREY)
add_footer(s, 3)

# =========================================================================
# SLIDE 4 - TOP 10 LEVIERS PRIORITAIRES
# =========================================================================
s = add_blank()
add_header(s, "3. Top 10 leviers commerciaux", "Les dispositifs a maitriser en priorite par les equipes commerciales")

top10 = [
    ("CEE BAT-EQ-127", "Luminaires LED tertiaire - finance 30-70% du relamping", "ELEC"),
    ("CEE BAT-TH-104 + Coup de pouce", "PAC en remplacement chaudieres - bonifie sortie fioul", "CVC"),
    ("CEE BAT-TH-116 + decret BACS", "GTB pilotage CVC - obligation > 290 kW (2025), > 70 kW (2027)", "ELEC/CVC"),
    ("Coup de pouce Renovation tertiaire", "Bouquet >= 30% gain - bonification jusqu'a x4", "TOUS"),
    ("Prime autoconso PV + Loi APER", "Ombrieres parkings > 1500 m2 (500 m2 en 07/2026) - marche captif", "ELEC"),
    ("Programme ADVENIR", "Bornes IRVE copro / entreprises - jusqu'a 1660 EUR + 3000 EUR pre-equipement", "ELEC"),
    ("Decret Tertiaire (DEET)", "Obligation -40% en 2030 - porte d'entree commerciale n 1", "TOUS"),
    ("Fonds Chaleur ADEME", "Jusqu'a 65% des couts EnR - geothermie, biomasse, chaleur fatale", "CVC"),
    ("Diag Decarbon'Action Bpifrance", "Diagnostic PME 4000 EUR forfait - genere le pipeline travaux", "TOUS"),
    ("MaPrimeRenov Copropriete", "Marche copro IDF - chauffage collectif, parties communes, IRVE", "CVC/ELEC"),
]

sy = Inches(1.3)
rh = Inches(0.5)
for i, (nom, desc, secteur) in enumerate(top10):
    y = sy + rh * i
    fill = LIGHT if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.5), y, Inches(12.3), rh - Inches(0.05), fill)
    add_text(s, Inches(0.6), y + Inches(0.08), Inches(0.6), Inches(0.4),
             f"{i+1:02d}", size=18, bold=True, color=ORANGE)
    add_text(s, Inches(1.2), y + Inches(0.08), Inches(3.5), Inches(0.4),
             nom, size=12, bold=True, color=NAVY)
    add_text(s, Inches(4.7), y + Inches(0.08), Inches(6.5), Inches(0.4),
             desc, size=11, color=NAVY)
    add_text(s, Inches(11.3), y + Inches(0.08), Inches(1.4), Inches(0.4),
             secteur, size=10, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)

add_footer(s, 4)

# =========================================================================
# SLIDE 5 - FOCUS DECRET TERTIAIRE
# =========================================================================
s = add_blank()
add_header(s, "4. Focus Decret Tertiaire (DEET)", "Notre porte d'entree commerciale n 1")

add_rect(s, Inches(0.5), Inches(1.3), Inches(6.2), Inches(5.5), LIGHT)
add_text(s, Inches(0.7), Inches(1.5), Inches(5.8), Inches(0.5),
         "Qui est concerne ?", size=16, bold=True, color=NAVY)
add_text(s, Inches(0.7), Inches(2.0), Inches(5.8), Inches(3.0),
         "Tous batiments tertiaires >= 1000 m2.\n\n"
         "Tous usages : bureaux, commerces, sante, enseignement, hotellerie, "
         "restauration, logistique, services, sport, culture, justice...\n\n"
         "Proprietaires ET locataires - obligation de declaration sur OPERAT.",
         size=12, color=NAVY)

add_text(s, Inches(0.7), Inches(4.7), Inches(5.8), Inches(0.5),
         "Trajectoire reglementaire", size=16, bold=True, color=NAVY)
add_text(s, Inches(0.7), Inches(5.2), Inches(5.8), Inches(1.5),
         "-40% en 2030 (base 2010-2019)\n"
         "-50% en 2040\n"
         "-60% en 2050\n"
         "Sanctions : 7500 EUR + name & shame",
         size=12, color=NAVY, bold=True)

# Pitch HUARD
add_rect(s, Inches(7.0), Inches(1.3), Inches(5.8), Inches(5.5), NAVY)
add_text(s, Inches(7.2), Inches(1.5), Inches(5.4), Inches(0.5),
         "Notre offre HUARD", size=16, bold=True, color=ORANGE)

steps = [
    ("1. Audit gratuit OPERAT", "Verifier le statut declaratif du client et le delta vs objectif 2030."),
    ("2. Plan d'actions", "CVC + eclairage + GTB + EnR - chiffrage par lot."),
    ("3. Financement", "CEE bonifies + Coup de pouce + Fonds Chaleur + Region IDF - reste a charge minimal."),
    ("4. Realisation", "Travaux HUARD tous corps d'etat techniques."),
    ("5. Suivi via GTB", "Mesure & verification - lien direct avec contrat de maintenance HUARD."),
]

ty = Inches(2.0)
for title, desc in steps:
    add_text(s, Inches(7.2), ty, Inches(5.4), Inches(0.35),
             title, size=12, bold=True, color=WHITE)
    add_text(s, Inches(7.2), ty + Inches(0.3), Inches(5.4), Inches(0.55),
             desc, size=10, color=LIGHT)
    ty += Inches(0.9)

add_footer(s, 5)

# =========================================================================
# SLIDE 6 - FOCUS PRIME SOLAIRE 2026
# =========================================================================
s = add_blank()
add_header(s, "5. Focus prime solaire 2026", "La nouvelle dynamique PV tertiaire + Loi APER")

add_text(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
         "Trois leviers cumulables :", size=16, bold=True, color=NAVY)

pv_blocks = [
    ("Prime a l'autoconsommation (arrete S26)",
     ORANGE,
     "Versee en 5 ans (1/5 par an). Indicatif T2 2026 :\n"
     "- <= 3 kWc : ~80 EUR/kWc\n"
     "- 3-9 kWc : ~80 EUR/kWc\n"
     "- 9-36 kWc : ~190 EUR/kWc\n"
     "- 36-100 kWc : ~100 EUR/kWc\n\n"
     "Pour 100 a 500 kWc : modalites specifiques."),
    ("Tarif d'achat EDF OA (surplus)",
     GREEN,
     "Contrat 20 ans, tarif EUR/kWh fixe.\n"
     "Securise le ROI du projet.\n\n"
     "Cumul prime + tarif possible jusqu'a 500 kWc.\n\n"
     "Au-dela : passer par appel d'offres CRE."),
    ("Loi APER - obligation",
     RED,
     "Parking > 1500 m2 : obligation\n"
     "(elargi a 500 m2 en juillet 2026)\n\n"
     "Batiments commerciaux/industriels > 500 m2 : integration PV obligatoire\n\n"
     "Sanction jusqu'a 40 000 EUR/an."),
]

x0 = Inches(0.5)
bw = Inches(4.0)
bh = Inches(4.5)
gap = Inches(0.2)
y0 = Inches(2.0)

for i, (title, color, body) in enumerate(pv_blocks):
    sx = x0 + (bw + gap) * i
    add_rect(s, sx, y0, bw, Inches(0.6), color)
    add_text(s, sx + Inches(0.2), y0 + Inches(0.1), bw - Inches(0.4), Inches(0.4),
             title, size=13, bold=True, color=WHITE)
    add_rect(s, sx, y0 + Inches(0.6), bw, bh - Inches(0.6), LIGHT)
    add_text(s, sx + Inches(0.3), y0 + Inches(0.85), bw - Inches(0.6), bh - Inches(0.95),
             body, size=11, color=NAVY)

add_text(s, Inches(0.5), Inches(6.85), Inches(12), Inches(0.4),
         "CIBLE PROSPECTION : hypermarches, zones commerciales, centres logistiques, sites industriels, sieges sociaux avec grand parking.",
         size=11, bold=True, color=NAVY)
add_footer(s, 6)

# =========================================================================
# SLIDE 7 - FOCUS IRVE
# =========================================================================
s = add_blank()
add_header(s, "6. Focus IRVE - bornes de recharge", "Programme ADVENIR + obligation LOM")

add_text(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
         "Plafonds ADVENIR 2026 (indicatifs)", size=16, bold=True, color=NAVY)

advenir_rows = [
    ("Particulier en immeuble collectif", "Jusqu'a 960 EUR HT / point de charge"),
    ("Copropriete - infrastructure collective", "Jusqu'a 1660 EUR/point + 3000 EUR pre-equipement"),
    ("Parking entreprise (salaries / flotte)", "Jusqu'a 1700 EUR HT / point"),
    ("Parking ouvert au public (standard)", "Jusqu'a 2700 EUR / point"),
    ("Bornes tres haute puissance (>140 kW)", "Jusqu'a 15 000 EUR / point"),
    ("Voirie - acteurs publics", "Sur dossier - jusqu'a 50% HT"),
]

ry = Inches(1.9)
for i, (cible, plafond) in enumerate(advenir_rows):
    add_rect(s, Inches(0.5), ry, Inches(12.3), Inches(0.55), LIGHT if i % 2 == 0 else WHITE)
    add_text(s, Inches(0.7), ry + Inches(0.12), Inches(6), Inches(0.4),
             cible, size=12, bold=True, color=NAVY)
    add_text(s, Inches(6.7), ry + Inches(0.12), Inches(6), Inches(0.4),
             plafond, size=12, color=ORANGE, bold=True)
    ry += Inches(0.55)

add_rect(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.4), NAVY)
add_text(s, Inches(0.7), Inches(5.65), Inches(12), Inches(0.4),
         "PREREQUIS HUARD", size=14, bold=True, color=ORANGE)
add_text(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.9),
         "1. Labellisation ADVENIR (formation Avere-France) - obligatoire pour declencher l'aide.\n"
         "2. Qualification IRVE Qualifelec niveau 2 ou 3 selon puissance.\n"
         "3. Engagement client avant signature du devis (clause ADVENIR a inclure dans tous nos devis IRVE).",
         size=11, color=LIGHT)

add_footer(s, 7)

# =========================================================================
# SLIDE 8 - CEE par metier
# =========================================================================
s = add_blank()
add_header(s, "7. CEE - fiches cles par metier HUARD", "Recensement des fiches a fort effet de levier")

table_data = [
    ("ELECTRICITE", BLUE,
     ["BAT-EQ-127 - Luminaires LED tertiaire",
      "BAT-EQ-133 - Gestion eclairage (DALI/KNX)",
      "BAT-EQ-117 - Variation puissance eclairage",
      "BAT-TH-116 - GTB classe A/B",
      "IND-UT-103 - Variateur vitesse moteurs"]),
    ("CVC", GREEN,
     ["BAT-TH-104 - PAC air/eau, eau/eau",
      "BAT-TH-127 - Raccordement reseau chaleur EnR",
      "BAT-TH-125/155 - VMC double flux HE",
      "BAT-TH-146 - Calorifugeage points singuliers",
      "BAT-TH-141 - Regulation programmation"]),
    ("MAINTENANCE", ORANGE,
     ["BAT-TH-146 + BAT-TH-145 - Calorifugeage (P2/P3)",
      "BAT-TH-116 - GTB exploitation",
      "BAT-TH-141 - Programmation horaire",
      "IND-UT-103 - VEV moteurs",
      "Coup de pouce 'Pilotage connecte chauffage'"]),
]

x0 = Inches(0.5)
bw = Inches(4.1)
bh = Inches(5.2)
gap = Inches(0.15)
y0 = Inches(1.4)

for i, (cat, color, items) in enumerate(table_data):
    sx = x0 + (bw + gap) * i
    add_rect(s, sx, y0, bw, Inches(0.6), color)
    add_text(s, sx + Inches(0.2), y0 + Inches(0.1), bw - Inches(0.4), Inches(0.4),
             cat, size=15, bold=True, color=WHITE)
    add_rect(s, sx, y0 + Inches(0.6), bw, bh - Inches(0.6), LIGHT)
    iy = y0 + Inches(0.85)
    for itm in items:
        add_rect(s, sx + Inches(0.25), iy + Inches(0.15), Inches(0.08), Inches(0.08), color)
        add_text(s, sx + Inches(0.45), iy, bw - Inches(0.7), Inches(0.65),
                 itm, size=11, color=NAVY)
        iy += Inches(0.85)

add_text(s, Inches(0.5), Inches(6.85), Inches(12), Inches(0.4),
         "Telecoms/Informatique : actions specifiques datacenter (refroidissement, free cooling) + Diag Decarbon'Action Bpifrance.",
         size=11, color=GREY)
add_footer(s, 8)

# =========================================================================
# SLIDE 9 - Eligibilite clients
# =========================================================================
s = add_blank()
add_header(s, "8. Eligibilite par profil client", "Qui beneficie de quoi - grille express")

headers = ["Profil client", "CEE", "Coup de pouce", "Fonds Chaleur", "Bpifrance", "Region IDF", "ADVENIR", "MPR / Eco-PTZ"]
rows = [
    ["PME tertiaire < 1000 m2", "X", "X", "X (EnR)", "X", "X", "X", "-"],
    ["Tertiaire 1000-5000 m2 (DEET)", "X", "X", "X (EnR)", "X", "X", "X", "-"],
    ["Grand tertiaire > 5000 m2", "X", "X", "X (EnR)", "Pret Vert", "X", "X", "-"],
    ["Industrie", "X (IND)", "(selon)", "X", "X", "X", "X", "-"],
    ["Collectivite (commune/MGP)", "X (BAT)", "(selon)", "X", "Banque Territoires", "X", "X", "-"],
    ["Copropriete", "X", "X", "X (EnR)", "Eco-PTZ", "X (POSIT'IF)", "X", "X (MPR Copro)"],
    ["Particulier", "X (BAR)", "X", "(limite)", "-", "(selon)", "X (960 EUR)", "X (MPR + Eco-PTZ)"],
    ["Hotel / sante / education", "X", "X", "X (ECS solaire, biomasse)", "X", "X", "X", "-"],
    ["Logistique / entrepots", "X", "X", "(selon)", "X", "X", "X", "-"],
]

ncols = len(headers)
col_widths = [Inches(2.6)] + [Inches(1.34)] * (ncols - 1)
total_w = sum(col_widths, Inches(0))
start_x = Inches(0.5)
y = Inches(1.4)
rh = Inches(0.42)

# Header row
x = start_x
for ci, h in enumerate(headers):
    add_rect(s, x, y, col_widths[ci], rh, NAVY)
    add_text(s, x + Inches(0.05), y + Inches(0.05), col_widths[ci] - Inches(0.1), rh - Inches(0.1),
             h, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += col_widths[ci]
y += rh

for ri, row in enumerate(rows):
    x = start_x
    for ci, val in enumerate(row):
        fill = LIGHT if ri % 2 == 0 else WHITE
        # color code for cells X / -
        if ci > 0:
            if val.startswith("X"):
                fill = RGBColor(0xD9, 0xF0, 0xE3)
            elif val == "-":
                fill = RGBColor(0xF3, 0xC7, 0xC2)
            elif val.startswith("("):
                fill = RGBColor(0xFF, 0xE9, 0xC4)
        add_rect(s, x, y, col_widths[ci], rh, fill)
        col_text = NAVY if ci == 0 else NAVY
        bold = (ci == 0)
        add_text(s, x + Inches(0.05), y + Inches(0.04), col_widths[ci] - Inches(0.1), rh - Inches(0.08),
                 val, size=10, color=col_text, bold=bold,
                 align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += col_widths[ci]
    y += rh

add_text(s, Inches(0.5), Inches(6.9), Inches(12), Inches(0.4),
         "Legende : X = eligible / (selon) = sous conditions ou selon AAP en cours / - = non eligible directement",
         size=10, color=GREY)
add_footer(s, 9)

# =========================================================================
# SLIDE 10 - Cumul des aides
# =========================================================================
s = add_blank()
add_header(s, "9. Cumul des aides", "Regles d'or pour eviter les pieges")

cumul = [
    ("Cumul CEE + Fonds Chaleur", GREEN,
     "POSSIBLE mais : ne pas subventionner deux fois le meme equipement. ADEME verifie. Bien isoler les lots."),
    ("Cumul CEE + MaPrimeRenov / MPR Copro", GREEN,
     "POSSIBLE. Logique : CEE deduit du devis, MPR verse au client. Decret MPR autorise."),
    ("Cumul CEE + ADVENIR", RED,
     "NON CUMULABLE. ADVENIR est finance via les CEE - on choisit l'un OU l'autre par point de charge."),
    ("Cumul CEE + Bpifrance Pret Vert / PEE", GREEN,
     "POSSIBLE. Le pret finance le reste a charge apres aides directes."),
    ("Cumul CEE + Region IDF", GREEN,
     "POSSIBLE sous regime de minimis (200 000 EUR sur 3 ans par entreprise)."),
    ("Cumul Coup de pouce + CEE classique", GREEN,
     "Le Coup de pouce REMPLACE la fiche CEE classique sur la meme operation - bonification appliquee."),
    ("Cumul TVA 5,5% / 10%", GREEN,
     "CUMULABLE avec toutes les aides - attestation client systematique."),
]

y = Inches(1.4)
for i, (rule, color, desc) in enumerate(cumul):
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.7), LIGHT if i % 2 == 0 else WHITE)
    add_rect(s, Inches(0.5), y, Inches(0.15), Inches(0.7), color)
    add_text(s, Inches(0.8), y + Inches(0.1), Inches(4.0), Inches(0.5),
             rule, size=12, bold=True, color=NAVY)
    add_text(s, Inches(5.0), y + Inches(0.1), Inches(7.7), Inches(0.55),
             desc, size=11, color=NAVY)
    y += Inches(0.75)

add_footer(s, 10)

# =========================================================================
# SLIDE 11 - PITCH 1 : DECRET TERTIAIRE
# =========================================================================
s = add_blank()
add_header(s, "10. Pitch 1 - Decret Tertiaire", "Le pitch a deployer chez tous nos clients > 1000 m2")

add_rect(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.3), NAVY)
add_text(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.0),
         "« Vous etes oblige par le decret tertiaire d'atteindre -40% de conso en 2030. "
         "Si vous ne deposez pas la declaration OPERAT, c'est 7500 EUR d'amende plus le name & shame. "
         "Nous proposons un parcours simple : audit -> plan d'actions -> travaux finances jusqu'a 70% par CEE + ADEME. »",
         size=14, color=WHITE)

add_text(s, Inches(0.5), Inches(2.9), Inches(6), Inches(0.5),
         "Cible commerciale", size=14, bold=True, color=NAVY)
add_text(s, Inches(0.7), Inches(3.4), Inches(5.8), Inches(3.0),
         "- Bureaux et commerces > 1000 m2\n"
         "- Hopitaux, ehpad, cliniques\n"
         "- Hotels et residences services\n"
         "- Etablissements scolaires prives et publics\n"
         "- Centres commerciaux\n"
         "- Logistique et entrepots tertiaires\n"
         "- Sieges sociaux groupes",
         size=12, color=NAVY)

add_text(s, Inches(7.0), Inches(2.9), Inches(6), Inches(0.5),
         "Outils a sortir", size=14, bold=True, color=NAVY)
add_text(s, Inches(7.2), Inches(3.4), Inches(5.8), Inches(3.0),
         "- One-pager DEET (a creer en interne)\n"
         "- Acces OPERAT (donnees publiques)\n"
         "- Modele de bilan d'audit en 1 page\n"
         "- Catalogue chiffrages CEE bonifies\n"
         "- Liste des partenaires BE qualifies\n"
         "- Contact mandataire CEE + ADEME IDF",
         size=12, color=NAVY)

add_footer(s, 11)

# =========================================================================
# SLIDE 12 - PITCH 2 : COPRO IDF
# =========================================================================
s = add_blank()
add_header(s, "11. Pitch 2 - Copropriete IDF", "Le marche structurant pour HUARD en Ile-de-France")

add_rect(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.3), NAVY)
add_text(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.0),
         "« Votre copropriete chauffe au fioul/gaz collectif ? Le coup de pouce CEE bonifie x4 le remplacement par PAC ou raccordement reseau de chaleur. "
         "Plus MaPrimeRenov Copropriete jusqu'a 25%, plus Energies POSIT'IF qui peut financer la totalite via tiers-financement. "
         "Resultat : projet vote en AG sans appel de fonds disproportionne. »",
         size=13, color=WHITE)

add_text(s, Inches(0.5), Inches(2.9), Inches(6.0), Inches(0.5),
         "Aides cumulables sur copro", size=14, bold=True, color=NAVY)
aides = [
    ("MaPrimeRenov Copropriete (Anah)", "Jusqu'a 25% du HT + bonus sortie passoire"),
    ("CEE + Coup de pouce", "Bonification x2-4 sur sortie fioul"),
    ("Eco-PTZ collectif", "Jusqu'a 50 000 EUR / logement (0%)"),
    ("Energies POSIT'IF (Region IDF)", "Tiers-financement integral"),
    ("Region IDF + Departements", "Subventions complementaires"),
    ("ADVENIR (IRVE collective)", "Jusqu'a 1660 EUR/point + 3000 EUR pre-eq."),
]
y = Inches(3.5)
for nom, mt in aides:
    add_rect(s, Inches(0.5), y, Inches(6.2), Inches(0.45), LIGHT)
    add_text(s, Inches(0.7), y + Inches(0.08), Inches(3.5), Inches(0.35),
             nom, size=11, bold=True, color=NAVY)
    add_text(s, Inches(4.2), y + Inches(0.08), Inches(2.5), Inches(0.35),
             mt, size=10, color=NAVY)
    y += Inches(0.5)

add_text(s, Inches(7.0), Inches(2.9), Inches(6.0), Inches(0.5),
         "Strategie HUARD copro", size=14, bold=True, color=NAVY)
add_text(s, Inches(7.2), Inches(3.5), Inches(5.8), Inches(3.5),
         "1. Cartographier les syndics partenaires IDF\n\n"
         "2. Proposer des reunions d'information dediees AG\n\n"
         "3. Outils : 'one-pager financement copro' par typologie\n\n"
         "4. Partenariat Energies POSIT'IF et MAR (Mon Accompagnateur Renov)\n\n"
         "5. Offre tout-en-un : chauffage + ECS + IRVE + parties communes",
         size=11, color=NAVY)

add_footer(s, 12)

# =========================================================================
# SLIDE 13 - PITCH 3 : LOI APER PARKINGS
# =========================================================================
s = add_blank()
add_header(s, "12. Pitch 3 - Loi APER & ombrieres", "Marche captif - 24 mois pour se positionner")

add_rect(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.3), NAVY)
add_text(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.0),
         "« Tout parking > 1500 m2 doit etre couvert d'ombrieres PV. En juillet 2026, le seuil descend a 500 m2. "
         "Sanction : jusqu'a 40 000 EUR/an non couverte. Sur 1000 m2 de parking, c'est ~100 kWc PV, "
         "soit 13-15k EUR/an de revente surplus + prime autoconso. Nous montons le projet cle en main. »",
         size=13, color=WHITE)

add_text(s, Inches(0.5), Inches(2.9), Inches(12), Inches(0.5),
         "Cibles a prospecter en priorite sur l'IDF", size=14, bold=True, color=NAVY)

cibles_pv = [
    ("Hypermarches & GMS", "Auchan, Carrefour, Leclerc, Casino"),
    ("Centres commerciaux & retail parks", "Westfield, Klepierre, Hammerson, Mercialys"),
    ("Logistique & messagerie", "Prologis, GLP, Goodman, Argan, Geodis, DHL"),
    ("Concessions automobiles", "Stellantis, Renault, BMW, Mercedes (parkings)"),
    ("Zones industrielles", "ZI Trappes, Massy, Senart, Roissy, Cergy"),
    ("Sieges et campus tertiaires", "Saclay, La Defense, Velizy, Rueil"),
]
y = Inches(3.5)
for i, (cat, exemples) in enumerate(cibles_pv):
    col = i % 2
    row_i = i // 2
    x = Inches(0.5) + Inches(6.4) * col
    yy = y + Inches(0.95) * row_i
    add_rect(s, x, yy, Inches(6.1), Inches(0.85), LIGHT)
    add_text(s, x + Inches(0.2), yy + Inches(0.1), Inches(5.7), Inches(0.4),
             cat, size=12, bold=True, color=NAVY)
    add_text(s, x + Inches(0.2), yy + Inches(0.42), Inches(5.7), Inches(0.4),
             exemples, size=10, color=GREY)

add_footer(s, 13)

# =========================================================================
# SLIDE 14 - PLAN D'ACTIONS COMMERCIAL
# =========================================================================
s = add_blank()
add_header(s, "13. Plan d'actions commercial 90 jours", "Roadmap d'activation interne et externe")

actions = [
    ("J0-J15", "Outillage", ORANGE,
     "- Diffusion du classeur Excel + PPT aux equipes\n"
     "- Formation 2h : top 10 leviers + cas concrets\n"
     "- Creation one-pagers par metier (DEET, BACS, APER, IRVE)"),
    ("J15-J45", "Activation interne", BLUE,
     "- Selection mandataire CEE partenaire (negociation rachat kWh cumac)\n"
     "- Confirmation labellisation ADVENIR + Qualifelec IRVE\n"
     "- Mise a jour devis type : mentions CEE/ADVENIR/TVA"),
    ("J45-J75", "Detection clients", GREEN,
     "- Scan portefeuille : clients > 1000 m2 tertiaire (DEET) et > 290 kW CVC (BACS)\n"
     "- Liste prospect ombrieres : parkings > 500 m2\n"
     "- Liste copros IDF avec chauffage collectif fioul/gaz"),
    ("J75-J90", "Mise en pression commerciale", RED,
     "- Campagne de RDV 'Audit aides energie' (gratuit, 30 min)\n"
     "- KPIs : nombre d'audits realises, devis emis avec aides, taux signature\n"
     "- Reporting mensuel COMEX"),
]

y = Inches(1.35)
for periode, theme, color, contenu in actions:
    add_rect(s, Inches(0.5), y, Inches(1.6), Inches(1.3), color)
    add_text(s, Inches(0.6), y + Inches(0.2), Inches(1.4), Inches(0.4),
             periode, size=14, bold=True, color=WHITE)
    add_text(s, Inches(0.6), y + Inches(0.65), Inches(1.4), Inches(0.5),
             theme, size=11, color=WHITE)
    add_rect(s, Inches(2.1), y, Inches(10.7), Inches(1.3), LIGHT)
    add_text(s, Inches(2.3), y + Inches(0.15), Inches(10.4), Inches(1.1),
             contenu, size=11, color=NAVY)
    y += Inches(1.42)

add_footer(s, 14)

# =========================================================================
# SLIDE 15 - SOURCES
# =========================================================================
s = add_blank()
add_header(s, "14. Sources officielles", "A garder en favoris - documentation et veille")

sources_list = [
    ("CEE - fiches standardisees", "ecologie.gouv.fr/operations-standardisees"),
    ("Coups de pouce CEE", "ecologie.gouv.fr/coup-pouce-economies-denergie"),
    ("OPERAT - Decret Tertiaire", "operat.ademe.fr"),
    ("Decret BACS", "ecologie.gouv.fr/decret-bacs"),
    ("ADEME - Fonds Chaleur", "fondschaleur.ademe.fr"),
    ("ADEME - Agir transition (Tremplin, Diag)", "agirpourlatransition.ademe.fr"),
    ("Bpifrance - offres TEE", "bpifrance.fr/catalogue-offres/transition-ecologique-et-energetique"),
    ("Programme ADVENIR (IRVE)", "advenir.mobi"),
    ("Programme ACTEE (collectivites)", "programme-actee.fr"),
    ("MaPrimeRenov / France Renov", "france-renov.gouv.fr"),
    ("EDF Obligation d'achat solaire", "edf-oa.fr"),
    ("Photovoltaique.info (Hespul)", "photovoltaique.info"),
    ("Region IDF - aides TEE", "iledefrance.fr/aides-services/transition-energetique"),
    ("Energies POSIT'IF", "energiespositif.fr"),
    ("ADEME Ile-de-France", "ile-de-france.ademe.fr"),
    ("Annuaires RGE - Qualibat, Qualifelec, QualiPV", "qualibat.com / qualifelec.fr / qualit-enr.org"),
]

cols = 2
cw = Inches(6.1)
ch = Inches(0.4)
gap_x = Inches(0.2)
gap_y = Inches(0.1)
x0 = Inches(0.5)
y0 = Inches(1.4)

for i, (theme, url) in enumerate(sources_list):
    col = i % cols
    row_i = i // cols
    x = x0 + (cw + gap_x) * col
    yy = y0 + (ch + gap_y) * row_i
    add_rect(s, x, yy, cw, ch, LIGHT if i % 2 == 0 else WHITE)
    add_text(s, x + Inches(0.15), yy + Inches(0.07), Inches(3.0), Inches(0.3),
             theme, size=11, bold=True, color=NAVY)
    add_text(s, x + Inches(3.2), yy + Inches(0.07), Inches(2.8), Inches(0.3),
             url, size=10, color=BLUE)

add_footer(s, 15)

# =========================================================================
# SLIDE 16 - PLAN ADMINISTRATIF HUARD - PIPELINE 6 PHASES
# =========================================================================
s = add_blank()
add_header(s, "15. Plan administratif HUARD", "Pipeline de montage des aides client - 6 phases / de J0 a J+330")

phases_ppt = [
    ("J0 - J+2", "DETECTION", ORANGE,
     "Commercial qualifie\nGrille eligibilite\nOPERAT si DEET"),
    ("J+3 - J+10", "PRE-ETUDE", RED,
     "Chiffrage CEE / ADVENIR\nNote opportunite\nDevis pre-detaille"),
    ("J+10 - J+20", "DOSSIER", BLUE,
     "Mandat client\nPieces RGE / audit\nPhotos avant"),
    ("J+15 - J+25", "ENGAGEMENT", GREEN,
     "Depot plateforme\n(CEE/ADVENIR/MPR)\nPUIS signature devis"),
    ("J+25 - J+150", "TRAVAUX", RGBColor(0x82, 0xC3, 0x41),
     "Execution + tracabilite\nFactures detaillees\nPV pose + photos"),
    ("J+150 - J+330", "VERSEMENT", NAVY,
     "Pieces de cloture\nVersement aide\nSuivi 3-5 ans"),
]

bw = Inches(2.05)
bh = Inches(4.5)
gap = Inches(0.1)
x0 = Inches(0.4)
y0 = Inches(1.4)

for i, (delai, nom, color, contenu) in enumerate(phases_ppt):
    sx = x0 + (bw + gap) * i
    add_rect(s, sx, y0, bw, Inches(0.5), color)
    add_text(s, sx + Inches(0.1), y0 + Inches(0.1), bw - Inches(0.2), Inches(0.35),
             delai, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, sx, y0 + Inches(0.5), bw, Inches(0.7), color)
    add_text(s, sx + Inches(0.1), y0 + Inches(0.6), bw - Inches(0.2), Inches(0.5),
             f"PHASE {i+1}\n{nom}", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, sx, y0 + Inches(1.2), bw, bh - Inches(1.2), LIGHT)
    add_text(s, sx + Inches(0.15), y0 + Inches(1.35), bw - Inches(0.3), bh - Inches(1.4),
             contenu, size=10, color=NAVY)
    # Fleche
    if i < len(phases_ppt) - 1:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, sx + bw - Inches(0.05), y0 + Inches(2.0), Inches(0.15), Inches(0.25))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ORANGE
        arrow.line.fill.background()

add_rect(s, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.65), RED)
add_text(s, Inches(0.6), Inches(6.3), Inches(12.2), Inches(0.5),
         "POINT CRITIQUE : pour CEE et ADVENIR, l'engagement sur plateforme DOIT etre fait AVANT signature devis - sinon prime perdue.",
         size=12, bold=True, color=WHITE)
add_footer(s, 16)

# =========================================================================
# SLIDE 17 - RACI INTERNE HUARD
# =========================================================================
s = add_blank()
add_header(s, "16. RACI interne HUARD", "Qui fait quoi - matrice claire pour eviter les trous dans la raquette")

raci_headers = ["Action", "Commercial", "Charge d'affaires", "Service aides", "Conduite trvx", "Comptabilite"]
raci_data = [
    ("Detection / qualification client", "R", "C", "I", "I", "-"),
    ("Chiffrage devis + aides", "C", "R", "C", "C", "-"),
    ("Choix du dispositif optimal", "C", "C", "R", "-", "-"),
    ("Constitution dossier admin", "C", "C", "R", "I", "-"),
    ("Mandat client signature", "R", "C", "C", "-", "-"),
    ("Engagement plateformes", "I", "C", "R", "-", "-"),
    ("Signature devis (apres engagement)", "R", "C", "I", "-", "-"),
    ("Execution + tracabilite", "I", "C", "I", "R", "-"),
    ("Depot pieces cloture", "I", "C", "R", "C", "-"),
    ("Versement + facturation", "I", "I", "R", "-", "C"),
    ("Controle a posteriori 3-5 ans", "-", "-", "R", "C", "I"),
]

col_widths = [Inches(4.0)] + [Inches(1.7)] * 5
start_x = Inches(0.5)
y = Inches(1.3)
rh = Inches(0.36)

x = start_x
for ci, h in enumerate(raci_headers):
    add_rect(s, x, y, col_widths[ci], rh, NAVY)
    add_text(s, x + Inches(0.05), y + Inches(0.04), col_widths[ci] - Inches(0.1), rh - Inches(0.08),
             h, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += col_widths[ci]
y += rh

color_raci = {
    "R": (GREEN, WHITE), "C": (RGBColor(0xFF, 0xE9, 0xC4), NAVY),
    "I": (LIGHT, NAVY), "-": (WHITE, GREY)
}

for ri, row in enumerate(raci_data):
    x = start_x
    for ci, val in enumerate(row):
        if ci == 0:
            fill = LIGHT if ri % 2 == 0 else WHITE
            add_rect(s, x, y, col_widths[ci], rh, fill)
            add_text(s, x + Inches(0.1), y + Inches(0.04), col_widths[ci] - Inches(0.2), rh - Inches(0.08),
                     val, size=10, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        else:
            fc, tc = color_raci.get(val, (LIGHT, NAVY))
            add_rect(s, x, y, col_widths[ci], rh, fc)
            add_text(s, x + Inches(0.05), y + Inches(0.04), col_widths[ci] - Inches(0.1), rh - Inches(0.08),
                     val, size=11, bold=True, color=tc, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += col_widths[ci]
    y += rh

leg_y = y + Inches(0.2)
add_text(s, Inches(0.5), leg_y, Inches(12), Inches(0.4),
         "R = Realise   /   C = Consulte   /   I = Informe   /   - = Non implique",
         size=11, color=NAVY)

add_text(s, Inches(0.5), leg_y + Inches(0.5), Inches(12), Inches(0.4),
         "RECOMMANDATION : creer un poste 'Referent aides energie' (1 ETP) pour porter la phase 2 a 6 - amorti des le 1er trimestre sur le volume vise.",
         size=11, bold=True, color=ORANGE)
add_footer(s, 17)

# =========================================================================
# SLIDE 18 - MATRICE PROJET x AIDES
# =========================================================================
s = add_blank()
add_header(s, "17. Matrice projet x aides", "16 types de projets HUARD - toutes les aides cumulables")

projets_pt = [
    ("Relamping LED tertiaire", "ELEC", "CEE BAT-EQ-127 + 133", "30-70%"),
    ("Chaudiere fioul -> PAC copro", "CVC", "CEE bonifie x4 + MPR Copro + POSIT'IF", "40-80%"),
    ("PAC tertiaire >100 kW", "CVC", "CEE + Coup de pouce + Fonds Chaleur", "30-65%"),
    ("PV autoconso 9-100 kWc", "ELEC", "Prime S26 + tarif EDF OA 20 ans", "75-85% (capex)"),
    ("Ombrieres PV parking (APER)", "ELEC", "Prime S26 + tarif rachat + Region", "amorti 7-10 ans"),
    ("GTB classe A (BACS)", "ELEC/CVC", "CEE BAT-TH-116 + Coup de pouce pilotage", "50-80%"),
    ("IRVE parking entreprise", "ELEC", "ADVENIR jusqu'a 1700 EUR/point", "30-50%"),
    ("IRVE collective copro", "ELEC", "ADVENIR 1660 EUR/point + 3000 EUR pre-equip", "40-60%"),
    ("Raccordement reseau chaleur", "CVC", "CEE BAT-TH-127 + Fonds Chaleur + MPR Copro", "40-70%"),
    ("Bouquet renovation DEET", "MULTI", "Coup de pouce x4 + ADEME + Region + Bpifrance", "40-70%"),
    ("Calorifugeage chaufferie", "CVC/MAINT", "CEE BAT-TH-145 + 146", "80-100%"),
    ("VMC double flux", "CVC", "CEE BAT-TH-125/155 + ACTEE + AIRPARIF", "25-50%"),
    ("Variateurs moteurs industriels", "ELEC", "CEE IND-UT-103 + Tremplin + Diag Eco-Flux", "30-60%"),
    ("Recup chaleur fatale", "CVC", "CEE IND-UT-117/102 + Fonds Chaleur + France 2030", "40-75%"),
    ("Chaudiere -> PAC particulier", "CVC", "MPR + CEE Coup de pouce + Eco-PTZ + TVA 5,5%", "40-90%"),
    ("Datacenter free cooling", "INFO/CVC", "CEE + Fonds Chaleur (fatale) + France 2030", "25-50%"),
]

ncols_p = 4
col_w = [Inches(4.0), Inches(1.5), Inches(5.5), Inches(2.0)]
hdr = ["Projet HUARD", "Activite", "Aides mobilisables", "Couverture"]

x = Inches(0.4)
y = Inches(1.3)
rh = Inches(0.32)
for ci, h in enumerate(hdr):
    add_rect(s, x, y, col_w[ci], rh, NAVY)
    add_text(s, x + Inches(0.05), y + Inches(0.03), col_w[ci] - Inches(0.1), rh - Inches(0.05),
             h, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += col_w[ci]
y += rh

for ri, (proj, act, aides, cov) in enumerate(projets_pt):
    fill = LIGHT if ri % 2 == 0 else WHITE
    x = Inches(0.4)
    add_rect(s, x, y, col_w[0], rh, fill)
    add_text(s, x + Inches(0.1), y + Inches(0.03), col_w[0] - Inches(0.2), rh - Inches(0.06),
             proj, size=9, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    x += col_w[0]
    add_rect(s, x, y, col_w[1], rh, fill)
    add_text(s, x, y + Inches(0.03), col_w[1], rh - Inches(0.06),
             act, size=9, color=BLUE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    x += col_w[1]
    add_rect(s, x, y, col_w[2], rh, fill)
    add_text(s, x + Inches(0.1), y + Inches(0.03), col_w[2] - Inches(0.2), rh - Inches(0.06),
             aides, size=9, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    x += col_w[2]
    add_rect(s, x, y, col_w[3], rh, RGBColor(0xD9, 0xF0, 0xE3))
    add_text(s, x, y + Inches(0.03), col_w[3], rh - Inches(0.06),
             cov, size=9, color=GREEN, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y += rh

add_text(s, Inches(0.4), Inches(7.0), Inches(12), Inches(0.4),
         "Detail complet et regles de cumul : voir onglet '10 - Matrice projet x aides' du classeur Excel.",
         size=10, color=GREY)
add_footer(s, 18)

# =========================================================================
# SLIDE 19 - BASES DE CALCUL (formules cles)
# =========================================================================
s = add_blank()
add_header(s, "18. Bases de calcul", "Formules cles a memoriser pour chiffrer en 5 minutes")

formules = [
    ("CEE standard",
     "Prime = kWh_cumac x cours_marche / 1000",
     "Cours T2 2026 : ~7 a 10 EUR / MWh cumac",
     BLUE),
    ("CEE bonifie (Coup de pouce)",
     "Prime = kWh_cumac x mult (x2 a x4) x cours / 1000",
     "Sortie fioul = x4 - Sortie gaz = x2",
     ORANGE),
    ("Prime autoconso PV",
     "Prime = EUR_par_kWc x P_kWc",
     "Versee en 5 annuites - revision trimestrielle CRE",
     GREEN),
    ("Tarif EDF OA surplus",
     "Revenu/an = kWh_vendus x tarif (20 ans)",
     "Contrat fixe 20 ans signe a la mise en service",
     GREEN),
    ("ADVENIR IRVE",
     "Aide = MIN(% HT ; plafond EUR/point)",
     "Plafond varie selon parking : 960 a 15 000 EUR",
     ORANGE),
    ("Fonds Chaleur forfait",
     "Aide = MWh_EnR/an x forfait EUR/MWh",
     "Geothermie : 80 - Biomasse : 25-35 - Solaire thermique : 40-80",
     RGBColor(0x82, 0xC3, 0x41)),
    ("MaPrimeRenov Copro",
     "Aide = % gain x cout_HT_plafonne / logement",
     "Plafond 25 000 EUR HT/log + bonus sortie passoire",
     RED),
    ("TVA 5,5%",
     "Economie = HT x 14,5%",
     "Travaux residentiel > 2 ans - attestation client",
     BLUE),
    ("Bpifrance PEE",
     "Pret 10 a 500 k EUR - taux fixe ~3-4%",
     "Eligible si action figure dans le perimetre CEE",
     GREEN),
]

cols_f = 3
fw = Inches(4.2)
fh = Inches(1.55)
gx = Inches(0.15)
gy = Inches(0.15)
x0 = Inches(0.35)
y0 = Inches(1.3)

for i, (nom, formule, note, color) in enumerate(formules):
    col = i % cols_f
    rowf = i // cols_f
    sx = x0 + (fw + gx) * col
    sy = y0 + (fh + gy) * rowf
    add_rect(s, sx, sy, fw, Inches(0.4), color)
    add_text(s, sx + Inches(0.15), sy + Inches(0.06), fw - Inches(0.3), Inches(0.3),
             nom, size=12, bold=True, color=WHITE)
    add_rect(s, sx, sy + Inches(0.4), fw, fh - Inches(0.4), LIGHT)
    add_text(s, sx + Inches(0.15), sy + Inches(0.5), fw - Inches(0.3), Inches(0.4),
             formule, size=11, bold=True, color=NAVY)
    add_text(s, sx + Inches(0.15), sy + Inches(0.95), fw - Inches(0.3), Inches(0.5),
             note, size=9, color=GREY)

add_text(s, Inches(0.35), Inches(6.8), Inches(12.5), Inches(0.5),
         "Simulateurs Excel vivants disponibles dans l'onglet '12 - Simulateurs' : entrer ses valeurs, le calcul est automatique.",
         size=11, bold=True, color=ORANGE)
add_footer(s, 19)

# =========================================================================
# SLIDE 20 - EXEMPLE CHIFFRE 1 : PAC COPRO
# =========================================================================
s = add_blank()
add_header(s, "19. Exemple chiffre 1 - PAC copropriete", "Remplacement chaudiere fioul collective - copro 50 logements")

add_rect(s, Inches(0.5), Inches(1.3), Inches(5.8), Inches(5.5), NAVY)
add_text(s, Inches(0.7), Inches(1.5), Inches(5.4), Inches(0.4),
         "DONNEES DE BASE", size=14, bold=True, color=ORANGE)

data1 = [
    ("Nombre de logements", "50"),
    ("Cout total chantier HT", "380 000 EUR"),
    ("Gain energetique attendu", "50%"),
    ("Multiplicateur Coup de pouce", "x4 (sortie fioul)"),
    ("Plafond MPR Copro / logement", "25 000 EUR HT"),
    ("Taux MPR (gain >= 50%)", "35%"),
]
y = Inches(2.0)
for label, val in data1:
    add_text(s, Inches(0.7), y, Inches(3.6), Inches(0.3), label, size=11, color=LIGHT)
    add_text(s, Inches(4.3), y, Inches(2.0), Inches(0.3), val, size=11, bold=True, color=WHITE)
    y += Inches(0.4)

add_rect(s, Inches(6.5), Inches(1.3), Inches(6.3), Inches(5.5), LIGHT)
add_text(s, Inches(6.7), Inches(1.5), Inches(5.9), Inches(0.4),
         "CALCUL DES AIDES", size=14, bold=True, color=NAVY)

calc1 = [
    ("Prime CEE standard estimee", "35 000 EUR", NAVY),
    ("Prime CEE bonifiee (x4)", "140 000 EUR", GREEN),
    ("Aide MPR Copro (25 000 x 50 x 35%)", "+ 437 500 EUR HT (plafond)", GREEN),
    ("Sous-total aides directes", "= 437 500 + 140 000 = 577 500 EUR", BLUE),
    ("[Si Energies POSIT'IF]", "Tiers-financement du reste possible", ORANGE),
    ("", "", None),
    ("Reste a charge HT (sans POSIT'IF)", "= 380 000 - 380 000* = 0 EUR cash", GREEN),
    ("* aides plafonnees au cout reel", "Le surplus aide est ecrete", GREY),
    ("Cout par logement copro", "0 a 500 EUR / log (selon)", GREEN),
]
y = Inches(2.0)
for label, val, color in calc1:
    if color is None:
        y += Inches(0.2)
        continue
    add_text(s, Inches(6.7), y, Inches(4.0), Inches(0.3), label, size=11, color=NAVY)
    add_text(s, Inches(10.7), y, Inches(2.0), Inches(0.3), val, size=11, bold=True, color=color)
    y += Inches(0.4)

add_footer(s, 20)

# =========================================================================
# SLIDE 21 - EXEMPLE CHIFFRE 2 : PV + IRVE entreprise
# =========================================================================
s = add_blank()
add_header(s, "20. Exemple chiffre 2 - PV + IRVE", "PME logistique 1500 m2 parking - 100 kWc PV + 10 bornes 22 kW")

# Bloc PV
add_rect(s, Inches(0.5), Inches(1.3), Inches(6.1), Inches(2.7), GREEN)
add_text(s, Inches(0.7), Inches(1.4), Inches(5.7), Inches(0.4),
         "VOLET 1 - PV autoconsommation 100 kWc", size=13, bold=True, color=WHITE)

pv_data = [
    ("Cout total HT", "110 000 EUR"),
    ("Prime autoconso (5 ans)", "10 000 EUR (100 EUR x 100 kWc)"),
    ("Production annuelle", "110 000 kWh"),
    ("Revenu rachat 30% surplus", "~ 4 280 EUR / an x 20 ans"),
    ("Economie autoconso 70%", "~ 16 200 EUR / an"),
    ("Gain cumule 20 ans", "~ 420 000 EUR (net)"),
    ("ROI", "~ 5 a 6 ans"),
]
y = Inches(1.85)
for label, val in pv_data:
    add_text(s, Inches(0.7), y, Inches(3.0), Inches(0.25), label, size=10, color=LIGHT)
    add_text(s, Inches(3.7), y, Inches(2.8), Inches(0.25), val, size=10, bold=True, color=WHITE)
    y += Inches(0.28)

# Bloc IRVE
add_rect(s, Inches(6.8), Inches(1.3), Inches(6.0), Inches(2.7), ORANGE)
add_text(s, Inches(7.0), Inches(1.4), Inches(5.6), Inches(0.4),
         "VOLET 2 - IRVE 10 bornes 22 kW (parking salaries)", size=13, bold=True, color=WHITE)

irve_data = [
    ("Cout total HT (10 x 3 500 EUR)", "35 000 EUR"),
    ("Plafond ADVENIR", "1 700 EUR HT / point"),
    ("Aide ADVENIR (10 x MIN(50% ; plafond))", "17 000 EUR HT"),
    ("Reste a charge HT", "18 000 EUR"),
    ("Taux de couverture", "49%"),
    ("Suramortissement IS (40% du HT)", "~ 7 200 EUR economie fiscale"),
    ("Reste a charge net fiscal", "~ 10 800 EUR"),
]
y = Inches(1.85)
for label, val in irve_data:
    add_text(s, Inches(7.0), y, Inches(3.2), Inches(0.25), label, size=10, color=LIGHT)
    add_text(s, Inches(10.2), y, Inches(2.6), Inches(0.25), val, size=10, bold=True, color=WHITE)
    y += Inches(0.28)

# Bloc synthese
add_rect(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.7), NAVY)
add_text(s, Inches(0.7), Inches(4.3), Inches(11.9), Inches(0.4),
         "SYNTHESE OFFRE GROUPEE PV + IRVE", size=15, bold=True, color=ORANGE)

synthese = [
    ("Investissement total", "145 000 EUR HT"),
    ("Aides directes (prime PV + ADVENIR)", "27 000 EUR HT"),
    ("Reste a charge HT", "118 000 EUR HT"),
    ("Reste a charge net (suramortissement IRVE)", "110 800 EUR"),
    ("Gain cumule 20 ans PV (rachat + autoconso)", "+ 420 000 EUR"),
    ("Bilan net 20 ans", "+ 309 000 EUR pour le client"),
]
y = Inches(4.95)
for label, val in synthese:
    add_text(s, Inches(0.7), y, Inches(7.0), Inches(0.3), label, size=12, color=LIGHT)
    color_v = GREEN if "+" in val and "EUR" in val else WHITE
    add_text(s, Inches(7.7), y, Inches(5.0), Inches(0.3), val, size=12, bold=True, color=color_v)
    y += Inches(0.3)

add_footer(s, 21)

# =========================================================================
# SLIDE 22 - CLOTURE
# =========================================================================
s = add_blank()
add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, NAVY)
add_rect(s, Inches(0), Inches(3.4), prs.slide_width, Inches(0.08), ORANGE)

add_text(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.0),
         "Nos clients ne paient pas le prix fort.", size=32, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(3.6), Inches(11.5), Inches(1.0),
         "Ils paient le delta apres aides - et HUARD le sait avant eux.",
         size=24, color=ORANGE)

add_text(s, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.4),
         "Suite de ce document :", size=14, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.4),
         "1. Classeur Excel detaille - 12 onglets dont plan administratif, matrice projet x aides, bases de calcul, simulateurs vivants",
         size=12, color=LIGHT)
add_text(s, Inches(0.8), Inches(5.95), Inches(11.5), Inches(0.4),
         "2. Programme de formation interne 2h par equipe (commerciaux, charges d'affaires, conduite)",
         size=12, color=LIGHT)
add_text(s, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.4),
         "3. Recrutement d'un Referent Aides Energie (1 ETP) - amorti des le 1er trimestre",
         size=12, color=LIGHT)
add_text(s, Inches(0.8), Inches(6.65), Inches(11.5), Inches(0.4),
         "4. KPIs commerciaux a integrer au reporting mensuel COMEX",
         size=12, color=LIGHT)
add_text(s, Inches(0.8), Inches(7.05), Inches(11.5), Inches(0.4),
         "Direction commerciale - Groupe HUARD - 16 mai 2026", size=10, color=GREY)

prs.save(OUT)
print(f"OK -> {OUT}")
